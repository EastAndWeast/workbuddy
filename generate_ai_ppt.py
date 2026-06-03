#!/usr/bin/env python3
"""Generate AI development history PPT based on meeting notes."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── Color palette ──
BG_DARK = RGBColor(0x0F, 0x11, 0x17)       # deep navy black
BG_CARD = RGBColor(0x1A, 0x1D, 0x27)       # card background
ACCENT  = RGBColor(0x00, 0xD4, 0xAA)        # teal/cyan accent
ACCENT2 = RGBColor(0x6C, 0x5C, 0xF7)        # purple accent
ACCENT3 = RGBColor(0xFF, 0x6B, 0x6B)        # coral/red accent
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
GRAY    = RGBColor(0x9E, 0xA1, 0xAE)
LIGHT  = RGBColor(0xC8, 0xCC, 0xD8)
SUBTITLE_GRAY = RGBColor(0x6B, 0x70, 0x80)

prs = Presentation()
prs.slide_width = Inches(16)
prs.slide_height = Inches(9)

W = prs.slide_width
H = prs.slide_height

# ── Helpers ──

def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape_rect(slide, left, top, width, height, fill_color, border_color=None, corner_radius=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    if corner_radius:
        shape.adjustments[0] = corner_radius
    return shape

def add_text_box(slide, left, top, width, height, text, font_size=18, color=WHITE,
                 bold=False, alignment=PP_ALIGN.LEFT, font_name="Microsoft YaHei"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_multi_text(slide, left, top, width, height, lines, font_size=16, color=WHITE,
                   line_spacing=1.5, font_name="Microsoft YaHei"):
    """lines: list of (text, color, bold, font_size_override)"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, (text, c, b, fs) in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(fs if fs else font_size)
        p.font.color.rgb = c if c else color
        p.font.bold = b
        p.font.name = font_name
        p.space_after = Pt(font_size * 0.4)
    return txBox

def add_bullet_list(slide, left, top, width, height, items, font_size=15, color=WHITE,
                    bullet_color=ACCENT, font_name="Microsoft YaHei"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = font_name
        p.space_after = Pt(8)
        p.level = 0
        # Use bullet
        pPr = p._pPr
        if pPr is None:
            from pptx.oxml.ns import qn
            pPr = p._p.get_or_add_pPr()
        from pptx.oxml.ns import qn
        buNone = pPr.findall(qn('a:buNone'))
        for bn in buNone:
            pPr.remove(bn)
        buChar = pPr.find(qn('a:buChar'))
        if buChar is None:
            from lxml import etree
            buChar = etree.SubElement(pPr, qn('a:buChar'))
        buChar.set('char', '▸')
        buClr = pPr.find(qn('a:buClr'))
        if buClr is None:
            from lxml import etree
            buClr = etree.SubElement(pPr, qn('a:buClr'))
        srgbClr = buClr.find(qn('a:srgbClr'))
        if srgbClr is None:
            from lxml import etree
            srgbClr = etree.SubElement(buClr, qn('a:srgbClr'))
        srgbClr.set('val', '{:02X}{:02X}{:02X}'.format(bullet_color[0], bullet_color[1], bullet_color[2]) if isinstance(bullet_color, tuple) else
                     '{:02X}{:02X}{:02X}'.format(bullet_color.red if hasattr(bullet_color, 'red') else 0, bullet_color.green if hasattr(bullet_color, 'green') else 0xD4, bullet_color.blue if hasattr(bullet_color, 'blue') else 0xAA))
        # Fix: use proper RGBColor attributes
        from pptx.dml.color import RGBColor
        if isinstance(bullet_color, RGBColor):
            srgbClr.set('val', '{:02X}{:02X}{:02X}'.format(bullet_color[0] if hasattr(bullet_color, '__getitem__') else 0x00,
                     bullet_color[1] if hasattr(bullet_color, '__getitem__') else 0xD4,
                     bullet_color[2] if hasattr(bullet_color, '__getitem__') else 0xAA))
    return txBox

def add_icon_card(slide, left, top, width, height, icon_text, title, desc, accent_color=ACCENT):
    """Card with icon emoji, title and description."""
    card = add_shape_rect(slide, left, top, width, height, BG_CARD, corner_radius=0.05)
    # Icon circle
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, left + Inches(0.3), top + Inches(0.35), Inches(0.6), Inches(0.6))
    circle.fill.solid()
    circle.fill.fore_color.rgb = accent_color
    circle.line.fill.background()
    tf = circle.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = icon_text
    p.font.size = Pt(22)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].space_before = Pt(0)
    # Title
    add_text_box(slide, left + Inches(0.3), top + Inches(1.15), width - Inches(0.6), Inches(0.4),
                 title, font_size=16, color=WHITE, bold=True)
    # Desc
    add_text_box(slide, left + Inches(0.3), top + Inches(1.55), width - Inches(0.6), height - Inches(1.85),
                 desc, font_size=12, color=GRAY)

def add_accent_line(slide, left, top, width, color=ACCENT):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Pt(3))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()

def add_slide_number(slide, num, total):
    add_text_box(slide, W - Inches(1.2), H - Inches(0.5), Inches(1), Inches(0.4),
                 f"{num}/{total}", font_size=11, color=GRAY, alignment=PP_ALIGN.RIGHT)

TOTAL_SLIDES = 12

# ════════════════════════════════════════
# SLIDE 1 — Cover
# ════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_slide_bg(slide, BG_DARK)

# Decorative gradient bar at top
bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, Inches(0.08))
bar.fill.solid()
bar.fill.fore_color.rgb = ACCENT
bar.line.fill.background()

# Left accent bar
lbar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.5), Inches(2.5), Pt(4), Inches(3.5))
lbar.fill.solid()
lbar.fill.fore_color.rgb = ACCENT
lbar.line.fill.background()

# Title
add_text_box(slide, Inches(2.0), Inches(2.5), Inches(10), Inches(1.2),
             "AI 智能体与自动化工作流", font_size=44, color=WHITE, bold=True)

add_text_box(slide, Inches(2.0), Inches(3.8), Inches(10), Inches(0.8),
             "发展历程 · 技术变革 · 应用展望", font_size=24, color=ACCENT)

add_text_box(slide, Inches(2.0), Inches(5.0), Inches(10), Inches(0.5),
             "从 OpenCloud 到企业数字员工：AI 应用层创业全景", font_size=16, color=GRAY)

# Date & author
add_text_box(slide, Inches(2.0), Inches(7.2), Inches(6), Inches(0.4),
             "2026 年 6 月  ·  Bruce", font_size=14, color=SUBTITLE_GRAY)

# Decorative dots on right side
for i in range(5):
    dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(12.5 + i * 0.45), Inches(3.5), Inches(0.12), Inches(0.12))
    dot.fill.solid()
    dot.fill.fore_color.rgb = ACCENT if i == 0 else GRAY
    dot.line.fill.background()

add_slide_number(slide, 1, TOTAL_SLIDES)

# ════════════════════════════════════════
# SLIDE 2 — Agenda
# ════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, Inches(0.08))
bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT; bar.line.fill.background()

add_text_box(slide, Inches(1.0), Inches(0.6), Inches(6), Inches(0.7),
             "内容导览", font_size=32, color=WHITE, bold=True)
add_accent_line(slide, Inches(1.0), Inches(1.3), Inches(1.5))

agenda_items = [
    ("01", "AI 工具发展背景与市场趋势", "从 ChatGPT 到本地 AI 智能体的演进脉络"),
    ("02", "本地 AI 的技术突破", "系统级权限如何改变自动化范式"),
    ("03", "AI 作为新一代操作系统", "记忆系统、技能生态与多入口协同"),
    ("04", "行业冲击与就业变革", "数字员工时代的人力成本重塑"),
    ("05", "从确定性到不确定性", "系统设计范式的根本转变"),
    ("06", "自动化内容生产实践", "端到端工作流构建与部署方案"),
    ("07", "AI 企业级集成与未来展望", "多 Agent 协同与个人助理趋势"),
]

for i, (num, title, desc) in enumerate(agenda_items):
    y = Inches(1.8) + Inches(i * 0.85)
    # Number
    num_shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.2), y, Inches(0.65), Inches(0.55))
    num_shape.fill.solid(); num_shape.fill.fore_color.rgb = BG_CARD; num_shape.line.fill.background()
    tf = num_shape.text_frame
    p = tf.paragraphs[0]; p.text = num; p.font.size = Pt(16); p.font.color.rgb = ACCENT
    p.alignment = PP_ALIGN.CENTER; p.font.bold = True; p.font.name = "Microsoft YaHei"
    # Title
    add_text_box(slide, Inches(2.1), y, Inches(6), Inches(0.35), title, font_size=17, color=WHITE, bold=True)
    add_text_box(slide, Inches(2.1), y + Inches(0.32), Inches(8), Inches(0.3), desc, font_size=13, color=GRAY)

add_slide_number(slide, 2, TOTAL_SLIDES)

# ════════════════════════════════════════
# SLIDE 3 — AI Development Background
# ════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, Inches(0.08))
bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT; bar.line.fill.background()

add_text_box(slide, Inches(1.0), Inches(0.5), Inches(10), Inches(0.7),
             "AI 应用层创业的爆发", font_size=32, color=WHITE, bold=True)
add_text_box(slide, Inches(1.0), Inches(1.1), Inches(10), Inches(0.4),
             "2024 春节 OpenCloud 引爆市场，开启本地 AI 智能体时代", font_size=15, color=GRAY)
add_accent_line(slide, Inches(1.0), Inches(1.6), Inches(1.5))

# Timeline cards
events = [
    ("2024.02", "OpenCloud 诞生", "具备读取本地数据能力的 AI\n首次突破云端权限限制\n引发行业广泛关注", ACCENT),
    ("2024 H2", "开源浪潮涌起", "OpenCloud 宣布开源\n创始人加入 OpenAI\n推出 Code、Cowork 等产品", ACCENT2),
    ("2025", "大厂纷纷入局", "腾讯推出 WorkBody\n小米推出 Me Cloud\n类似产品近百款涌现", ACCENT3),
    ("2026", "企业级落地", "AI 成为底层能力支撑\n数字员工概念普及\n自动化工作流成为标配", RGBColor(0xFF, 0xB7, 0x4D)),
]

for i, (date, title, desc, color) in enumerate(events):
    x = Inches(1.0) + Inches(i * 3.5)
    y = Inches(2.2)
    card = add_shape_rect(slide, x, y, Inches(3.1), Inches(4.0), BG_CARD, corner_radius=0.04)
    # Top color bar
    top_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, Inches(3.1), Pt(4))
    top_bar.fill.solid(); top_bar.fill.fore_color.rgb = color; top_bar.line.fill.background()
    add_text_box(slide, x + Inches(0.3), y + Inches(0.3), Inches(2.5), Inches(0.4),
                 date, font_size=14, color=color, bold=True)
    add_text_box(slide, x + Inches(0.3), y + Inches(0.75), Inches(2.5), Inches(0.5),
                 title, font_size=20, color=WHITE, bold=True)
    add_text_box(slide, x + Inches(0.3), y + Inches(1.4), Inches(2.5), Inches(2.2),
                 desc, font_size=13, color=LIGHT)

# Arrow connectors
for i in range(3):
    x = Inches(4.1) + Inches(i * 3.5)
    add_text_box(slide, x, Inches(3.8), Inches(0.5), Inches(0.4),
                 "→", font_size=24, color=GRAY, alignment=PP_ALIGN.CENTER)

add_slide_number(slide, 3, TOTAL_SLIDES)

# ════════════════════════════════════════
# SLIDE 4 — Local AI Technical Breakthrough
# ════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, Inches(0.08))
bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT; bar.line.fill.background()

add_text_box(slide, Inches(1.0), Inches(0.5), Inches(10), Inches(0.7),
             "本地 AI 能力的技术变革", font_size=32, color=WHITE, bold=True)
add_text_box(slide, Inches(1.0), Inches(1.1), Inches(12), Inches(0.4),
             "系统级权限让 AI 突破 API 限制，实现对任意软件的自动化操控", font_size=15, color=GRAY)
add_accent_line(slide, Inches(1.0), Inches(1.6), Inches(1.5))

# Left column - before
left_card = add_shape_rect(slide, Inches(1.0), Inches(2.2), Inches(6.5), Inches(5.8), BG_CARD, corner_radius=0.04)
add_text_box(slide, Inches(1.5), Inches(2.5), Inches(5.5), Inches(0.5),
             "传统方式 — 受限的自动化", font_size=18, color=ACCENT3, bold=True)

old_items = [
    "依赖 OCR 识别界面元素，精度低、易出错",
    "浏览器爬虫受限于权限和反爬机制",
    "缺少 API 的封闭系统无法集成",
    "跨软件操作需要大量适配开发",
    "邮件、Excel 等本地程序无法远程控制",
]
add_bullet_list(slide, Inches(1.5), Inches(3.2), Inches(5.5), Inches(4.0), old_items,
                font_size=14, color=LIGHT, bullet_color=ACCENT3)

# Right column - after
right_card = add_shape_rect(slide, Inches(8.0), Inches(2.2), Inches(6.5), Inches(5.8), BG_CARD, corner_radius=0.04)
add_text_box(slide, Inches(8.5), Inches(2.5), Inches(5.5), Inches(0.5),
             "本地 AI — 突破限制", font_size=18, color=ACCENT, bold=True)

new_items = [
    "系统级权限直接读取任意程序数据",
    "模拟人工操作，无需 API 即可控制软件",
    "跨终端跨平台无缝衔接",
    "监控邮件并自动回复处理",
    "Figma 内容读取 → IDE 导入 → 反馈",
]
add_bullet_list(slide, Inches(8.5), Inches(3.2), Inches(5.5), Inches(4.0), new_items,
                font_size=14, color=LIGHT, bullet_color=ACCENT)

# Center arrow
add_text_box(slide, Inches(7.0), Inches(4.5), Inches(1.2), Inches(0.6),
             "▶", font_size=28, color=ACCENT, alignment=PP_ALIGN.CENTER)

add_slide_number(slide, 4, TOTAL_SLIDES)

# ════════════════════════════════════════
# SLIDE 5 — AI as Operating System
# ════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, Inches(0.08))
bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT; bar.line.fill.background()

add_text_box(slide, Inches(1.0), Inches(0.5), Inches(10), Inches(0.7),
             "AI 作为新一代操作系统", font_size=32, color=WHITE, bold=True)
add_text_box(slide, Inches(1.0), Inches(1.1), Inches(12), Inches(0.4),
             "记忆、技能与多入口生态构建完整的 AI 协作体系", font_size=15, color=GRAY)
add_accent_line(slide, Inches(1.0), Inches(1.6), Inches(1.5))

# Three pillars
pillars = [
    ("🧠", "记忆系统", "短期记忆（工作内存）\n长期记忆（文件/数据库）\n支持持续性工作流\n跨会话上下文保留", ACCENT),
    ("⚡", "技能生态", "插件平台（Skill Hub）\n100+ 领域专家能力\n第三方能力集成\n可扩展的技能架构", ACCENT2),
    ("🔗", "多入口协同", "微信、Telegram 控制入口\n飞书/Slack 深度集成\n多 Agent 协同编排\n个人助理 + 数字同事", ACCENT3),
]

for i, (icon, title, desc, color) in enumerate(pillars):
    x = Inches(1.0) + Inches(i * 4.8)
    y = Inches(2.3)
    card = add_shape_rect(slide, x, y, Inches(4.3), Inches(5.5), BG_CARD, corner_radius=0.04)
    # Color top
    top_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, Inches(4.3), Pt(4))
    top_bar.fill.solid(); top_bar.fill.fore_color.rgb = color; top_bar.line.fill.background()
    # Icon
    add_text_box(slide, x + Inches(0.4), y + Inches(0.4), Inches(1), Inches(0.6),
                 icon, font_size=32, color=color)
    # Title
    add_text_box(slide, x + Inches(0.4), y + Inches(1.2), Inches(3.5), Inches(0.5),
                 title, font_size=22, color=WHITE, bold=True)
    # Desc
    add_text_box(slide, x + Inches(0.4), y + Inches(1.9), Inches(3.5), Inches(3.0),
                 desc, font_size=14, color=LIGHT)

add_slide_number(slide, 5, TOTAL_SLIDES)

# ════════════════════════════════════════
# SLIDE 6 — Industry Impact
# ════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, Inches(0.08))
bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT3; bar.line.fill.background()

add_text_box(slide, Inches(1.0), Inches(0.5), Inches(10), Inches(0.7),
             "行业冲击与就业形势变化", font_size=32, color=WHITE, bold=True)
add_text_box(slide, Inches(1.0), Inches(1.1), Inches(12), Inches(0.4),
             "数字员工取代重复性工作，企业组织结构面临深度重构", font_size=15, color=GRAY)
add_accent_line(slide, Inches(1.0), Inches(1.6), Inches(1.5), ACCENT3)

# Key stat highlight
stat_card = add_shape_rect(slide, Inches(1.0), Inches(2.2), Inches(14), Inches(1.8), BG_CARD, corner_radius=0.04)
add_text_box(slide, Inches(2.0), Inches(2.5), Inches(12), Inches(0.5),
             "美团传闻 2026 年 6 月 30 日裁员 50%", font_size=22, color=ACCENT3, bold=True)
add_text_box(slide, Inches(2.0), Inches(3.1), Inches(12), Inches(0.6),
             "若业绩不降反升 → 其他公司将跟进 → 互联网行业连锁反应", font_size=16, color=LIGHT)

# Impact areas
areas = [
    ("自动化写作", "PPT、日报、周报\n新闻稿件、研究报告", ACCENT),
    ("流程自动化", "群聊总结 → 任务拆解\n飞书审批 → 自动跟进", ACCENT2),
    ("代码与开发", "需求分析 → 代码实现\n测试 → 部署全流程", ACCENT3),
    ("招聘筛选", "简历筛选 → 面试安排\n人才匹配评分", RGBColor(0xFF, 0xB7, 0x4D)),
]

for i, (title, desc, color) in enumerate(areas):
    x = Inches(1.0) + Inches(i * 3.6)
    y = Inches(4.4)
    card = add_shape_rect(slide, x, y, Inches(3.2), Inches(3.2), BG_CARD, corner_radius=0.04)
    add_text_box(slide, x + Inches(0.3), y + Inches(0.3), Inches(2.6), Inches(0.5),
                 title, font_size=17, color=color, bold=True)
    add_accent_line(slide, x + Inches(0.3), y + Inches(0.85), Inches(1.0), color)
    add_text_box(slide, x + Inches(0.3), y + Inches(1.1), Inches(2.6), Inches(1.8),
                 desc, font_size=13, color=LIGHT)

add_slide_number(slide, 6, TOTAL_SLIDES)

# ════════════════════════════════════════
# SLIDE 7 — Deterministic to Non-deterministic
# ════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, Inches(0.08))
bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT2; bar.line.fill.background()

add_text_box(slide, Inches(1.0), Inches(0.5), Inches(12), Inches(0.7),
             "从确定性系统到不确定性系统", font_size=32, color=WHITE, bold=True)
add_text_box(slide, Inches(1.0), Inches(1.1), Inches(12), Inches(0.4),
             "AI 带来范式转变：产品设计、测试方法与预期管理的根本调整", font_size=15, color=GRAY)
add_accent_line(slide, Inches(1.0), Inches(1.6), Inches(1.5), ACCENT2)

# Old paradigm
old_card = add_shape_rect(slide, Inches(1.0), Inches(2.3), Inches(6.5), Inches(5.5), BG_CARD, corner_radius=0.04)
add_text_box(slide, Inches(1.5), Inches(2.6), Inches(5.5), Inches(0.5),
             "传统范式 — 确定性", font_size=20, color=GRAY, bold=True)
add_accent_line(slide, Inches(1.5), Inches(3.15), Inches(1.0), GRAY)

old_items = [
    "输入 → 输出一致且可复现",
    "Excel 公式永远返回相同结果",
    "测试覆盖率 = 质量保证",
    "设计模式强调完整控制",
]
add_bullet_list(slide, Inches(1.5), Inches(3.4), Inches(5.5), Inches(3.5), old_items,
                font_size=14, color=LIGHT, bullet_color=GRAY)

# New paradigm
new_card = add_shape_rect(slide, Inches(8.0), Inches(2.3), Inches(6.5), Inches(5.5), BG_CARD, corner_radius=0.04)
add_text_box(slide, Inches(8.5), Inches(2.6), Inches(5.5), Inches(0.5),
             "AI 范式 — 不确定性", font_size=20, color=ACCENT2, bold=True)
add_accent_line(slide, Inches(8.5), Inches(3.15), Inches(1.0), ACCENT2)

new_items = [
    "相同提示词可能产生不同输出",
    "交易推荐基于概率而非保证",
    "AI 写稿每次结果存在差异",
    "框架内自由发挥，部分可控",
]
add_bullet_list(slide, Inches(8.5), Inches(3.4), Inches(5.5), Inches(3.5), new_items,
                font_size=14, color=LIGHT, bullet_color=ACCENT2)

# Center transform arrow
add_text_box(slide, Inches(7.0), Inches(4.5), Inches(1.2), Inches(0.6),
             "⇄", font_size=32, color=ACCENT2, alignment=PP_ALIGN.CENTER)

add_slide_number(slide, 7, TOTAL_SLIDES)

# ════════════════════════════════════════
# SLIDE 8 — Automated Content Production
# ════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, Inches(0.08))
bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT; bar.line.fill.background()

add_text_box(slide, Inches(1.0), Inches(0.5), Inches(10), Inches(0.7),
             "自动化内容生产系统", font_size=32, color=WHITE, bold=True)
add_text_box(slide, Inches(1.0), Inches(1.1), Inches(12), Inches(0.4),
             "端到端工作流：信息采集 → 内容整合 → 排版设计 → 审核发布", font_size=15, color=GRAY)
add_accent_line(slide, Inches(1.0), Inches(1.6), Inches(1.5))

# Pipeline steps
steps = [
    ("01", "信息采集", "多源新闻抓取\n定时搜索策略\n失败重试机制", ACCENT),
    ("02", "内容整合", "AI 自动写稿\n模板化生成\n质量评分校验", ACCENT2),
    ("03", "排版设计", "封面图自动生成\n多品牌模板适配\n视频/图文输出", ACCENT3),
    ("04", "审核发布", "飞书审批工作流\n自动推送通知\nWordPress 同步", RGBColor(0xFF, 0xB7, 0x4D)),
]

for i, (num, title, desc, color) in enumerate(steps):
    x = Inches(1.0) + Inches(i * 3.6)
    y = Inches(2.3)
    card = add_shape_rect(slide, x, y, Inches(3.2), Inches(3.8), BG_CARD, corner_radius=0.04)
    top_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, Inches(3.2), Pt(4))
    top_bar.fill.solid(); top_bar.fill.fore_color.rgb = color; top_bar.line.fill.background()
    # Step number circle
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(1.1), y + Inches(0.3), Inches(0.7), Inches(0.7))
    circle.fill.solid(); circle.fill.fore_color.rgb = color; circle.line.fill.background()
    tf = circle.text_frame; p = tf.paragraphs[0]; p.text = num
    p.font.size = Pt(20); p.font.color.rgb = WHITE; p.alignment = PP_ALIGN.CENTER; p.font.bold = True
    add_text_box(slide, x + Inches(0.3), y + Inches(1.2), Inches(2.6), Inches(0.5),
                 title, font_size=18, color=WHITE, bold=True)
    add_text_box(slide, x + Inches(0.3), y + Inches(1.8), Inches(2.6), Inches(1.5),
                 desc, font_size=13, color=LIGHT)

# Arrows between steps
for i in range(3):
    x = Inches(4.2) + Inches(i * 3.6)
    add_text_box(slide, x, Inches(3.6), Inches(0.5), Inches(0.4),
                 "→", font_size=24, color=GRAY, alignment=PP_ALIGN.CENTER)

# Bottom note
note_card = add_shape_rect(slide, Inches(1.0), Inches(6.5), Inches(14), Inches(1.8), BG_CARD, corner_radius=0.04)
add_text_box(slide, Inches(1.5), Inches(6.8), Inches(13), Inches(0.5),
             "实际部署方案", font_size=17, color=WHITE, bold=True)
add_text_box(slide, Inches(1.5), Inches(7.3), Inches(13), Inches(0.7),
             "开源方案（OpenCloud）：可控性强但调试复杂  |  商业方案（WorkBody）：易用但存在服务商风险\n推荐策略：商业方案快速验证 → 苹果设备迁移 → GitHub 同步保障可移植性", font_size=14, color=LIGHT)

add_slide_number(slide, 8, TOTAL_SLIDES)

# ════════════════════════════════════════
# SLIDE 9 — Enterprise AI Integration
# ════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, Inches(0.08))
bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT; bar.line.fill.background()

add_text_box(slide, Inches(1.0), Inches(0.5), Inches(10), Inches(0.7),
             "AI 在企业工作流中的深度集成", font_size=32, color=WHITE, bold=True)
add_text_box(slide, Inches(1.0), Inches(1.1), Inches(12), Inches(0.4),
             "多 Agent 协同、专家工作流、外部连接器，AI 逐步成为企业底层能力", font_size=15, color=GRAY)
add_accent_line(slide, Inches(1.0), Inches(1.6), Inches(1.5))

scenarios = [
    ("代码开发", "需求分析 → 代码实现 → 测试部署\n全流程 AI 辅助\nLSP 集成提供代码智能", ACCENT),
    ("会议管理", "会议纪要自动生成\n待办任务拆解分发\n飞书日历/任务同步", ACCENT2),
    ("品牌设计", "AI 生产品牌素材\n设计稿自动迭代\n多平台内容适配", ACCENT3),
    ("跨境支付", "多币种结算处理\n合规审核自动化\nRWA 资产代币化", RGBColor(0xFF, 0xB7, 0x4D)),
    ("招聘筛选", "简历智能匹配\n面试日程协调\n人才画像生成", RGBColor(0x00, 0xB4, 0xD8)),
    ("内容运营", "多语言内容生成\nSEO 优化建议\n数据驱动选题", RGBColor(0xE8, 0x79, 0xF9)),
]

for i, (title, desc, color) in enumerate(scenarios):
    col = i % 3
    row = i // 3
    x = Inches(1.0) + Inches(col * 4.8)
    y = Inches(2.2) + Inches(row * 3.1)
    card = add_shape_rect(slide, x, y, Inches(4.3), Inches(2.7), BG_CARD, corner_radius=0.04)
    top_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, Inches(4.3), Pt(4))
    top_bar.fill.solid(); top_bar.fill.fore_color.rgb = color; top_bar.line.fill.background()
    add_text_box(slide, x + Inches(0.3), y + Inches(0.3), Inches(3.7), Inches(0.5),
                 title, font_size=18, color=color, bold=True)
    add_text_box(slide, x + Inches(0.3), y + Inches(1.0), Inches(3.7), Inches(1.4),
                 desc, font_size=13, color=LIGHT)

add_slide_number(slide, 9, TOTAL_SLIDES)

# ════════════════════════════════════════
# SLIDE 10 — Cost & Monitoring
# ════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, Inches(0.08))
bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT; bar.line.fill.background()

add_text_box(slide, Inches(1.0), Inches(0.5), Inches(10), Inches(0.7),
             "运行监控与成本管理", font_size=32, color=WHITE, bold=True)
add_text_box(slide, Inches(1.0), Inches(1.1), Inches(12), Inches(0.4),
             "积分制消耗模型 + 多层兜底策略保障系统稳定运行", font_size=15, color=GRAY)
add_accent_line(slide, Inches(1.0), Inches(1.6), Inches(1.5))

# Left - Monitoring
left_card = add_shape_rect(slide, Inches(1.0), Inches(2.3), Inches(7.0), Inches(5.5), BG_CARD, corner_radius=0.04)
add_text_box(slide, Inches(1.5), Inches(2.6), Inches(6), Inches(0.5),
             "运行监控机制", font_size=20, color=ACCENT, bold=True)
add_accent_line(slide, Inches(1.5), Inches(3.15), Inches(1.0), ACCENT)

mon_items = [
    "定时任务调度 — 7:00 早报 / 8:30 速报 / 19:00 时评 / 20:30 深度",
    "发布状态实时回显，成功/失败一目了然",
    "失败兜底处理 — 轮询重试 + 竞态条件防护",
    "Token 过期自动检测与刷新机制",
    "飞书审批流 → 人工确认 → 自动发布",
]
add_bullet_list(slide, Inches(1.5), Inches(3.4), Inches(6), Inches(3.5), mon_items,
                font_size=14, color=LIGHT, bullet_color=ACCENT)

# Right - Cost
right_card = add_shape_rect(slide, Inches(8.5), Inches(2.3), Inches(6.5), Inches(5.5), BG_CARD, corner_radius=0.04)
add_text_box(slide, Inches(9.0), Inches(2.6), Inches(5.5), Inches(0.5),
             "成本优化策略", font_size=20, color=ACCENT2, bold=True)
add_accent_line(slide, Inches(9.0), Inches(3.15), Inches(1.0), ACCENT2)

cost_items = [
    "每日免费额度有限，需合理分配",
    "高峰期积分不足 → 考虑购买套餐",
    "可切换 DeepSeek 等高性价比模型",
    "封面图生成使用 API 批量处理",
    "配置化部署降低迁移和维护成本",
]
add_bullet_list(slide, Inches(9.0), Inches(3.4), Inches(5.5), Inches(3.5), cost_items,
                font_size=14, color=LIGHT, bullet_color=ACCENT2)

add_slide_number(slide, 10, TOTAL_SLIDES)

# ════════════════════════════════════════
# SLIDE 11 — Future Outlook
# ════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, Inches(0.08))
bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT; bar.line.fill.background()

add_text_box(slide, Inches(1.0), Inches(0.5), Inches(10), Inches(0.7),
             "未来展望与行动建议", font_size=32, color=WHITE, bold=True)
add_text_box(slide, Inches(1.0), Inches(1.1), Inches(12), Inches(0.4),
             "拥抱变化，提升不可替代的能力", font_size=15, color=GRAY)
add_accent_line(slide, Inches(1.0), Inches(1.6), Inches(1.5))

# Three outlook cards
outlooks = [
    ("趋势", [
        "软件将集成大模型 + 控制中心",
        "个人助理 → 企业数字员工",
        "更多 AI + 区块链融合产品",
        "自动化工作流成为企业标配",
    ], ACCENT),
    ("挑战", [
        "权限与数据安全需要平衡",
        "不确定型系统的测试与验收",
        "多平台适配的标准化需求",
        "AI 决策的可解释性与合规",
    ], ACCENT2),
    ("行动", [
        "学习 AI 工具提升个人竞争力",
        "将重复性工作交给自动化",
        "关注开源项目保持技术敏感",
        "构建企业级 AI 应用能力",
    ], ACCENT3),
]

for i, (title, items, color) in enumerate(outlooks):
    x = Inches(1.0) + Inches(i * 4.8)
    y = Inches(2.2)
    card = add_shape_rect(slide, x, y, Inches(4.3), Inches(5.5), BG_CARD, corner_radius=0.04)
    top_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, Inches(4.3), Pt(4))
    top_bar.fill.solid(); top_bar.fill.fore_color.rgb = color; top_bar.line.fill.background()
    add_text_box(slide, x + Inches(0.4), y + Inches(0.4), Inches(3.5), Inches(0.5),
                 title, font_size=22, color=color, bold=True)
    add_accent_line(slide, x + Inches(0.4), y + Inches(0.95), Inches(1.0), color)
    add_bullet_list(slide, x + Inches(0.4), y + Inches(1.2), Inches(3.5), Inches(3.5), items,
                    font_size=14, color=LIGHT, bullet_color=color)

add_slide_number(slide, 11, TOTAL_SLIDES)

# ════════════════════════════════════════
# SLIDE 12 — Thank You
# ════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, Inches(0.08))
bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT; bar.line.fill.background()

# Centered content
add_text_box(slide, Inches(3), Inches(2.5), Inches(10), Inches(1.5),
             "感谢聆听", font_size=52, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(3), Inches(4.2), Inches(10), Inches(0.8),
             "AI 正在重塑每一个行业", font_size=24, color=ACCENT, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(3), Inches(5.5), Inches(10), Inches(0.6),
             "拥抱变化  ·  持续学习  ·  提升不可替代性", font_size=16, color=GRAY, alignment=PP_ALIGN.CENTER)

# Decorative dots
for i in range(5):
    dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(7.0 + i * 0.45), Inches(6.5), Inches(0.12), Inches(0.12))
    dot.fill.solid()
    dot.fill.fore_color.rgb = ACCENT if i == 2 else GRAY
    dot.line.fill.background()

add_text_box(slide, Inches(3), Inches(7.5), Inches(10), Inches(0.4),
             "2026 年 6 月  ·  Bruce", font_size=14, color=SUBTITLE_GRAY, alignment=PP_ALIGN.CENTER)

add_slide_number(slide, 12, TOTAL_SLIDES)

# ── Save ──
output_path = "/Users/bruce/WorkBuddy/Claw/AI智能体与自动化工作流.pptx"
prs.save(output_path)
print(f"PPT saved to: {output_path}")
